"""
Adds new slides to the existing Adobe Assessment presentation:
- CI/CD Pipeline (GitHub Actions)
- AWS Architecture
- Security & Encryption
- Branch Protection & Git Workflow
- Scalability
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import copy

ADOBE_RED   = RGBColor(0xEB, 0x1C, 0x24)
DARK_GRAY   = RGBColor(0x3C, 0x3C, 0x3C)
LIGHT_GRAY  = RGBColor(0xF4, 0xF4, 0xF4)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
BLUE        = RGBColor(0x00, 0x6E, 0xC7)
GREEN       = RGBColor(0x10, 0x7C, 0x10)

INPUT  = "Adobe_Assessment_Presentation.pptx"
OUTPUT = "Adobe_Assessment_Presentation.pptx"


def add_slide(prs, layout_index=0):
    layout = prs.slide_layouts[layout_index]
    return prs.slides.add_slide(layout)


def clear_slide(slide):
    for shape in list(slide.shapes):
        sp = shape._element
        sp.getparent().remove(sp)


def add_rect(slide, left, top, width, height, fill_color, line_color=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, text, left, top, width, height,
                 font_size=12, bold=False, color=DARK_GRAY,
                 align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txBox


def add_header_bar(slide, title):
    """Red header bar matching Adobe brand."""
    add_rect(slide, 0, 0, 10, 1.1, ADOBE_RED)
    add_text_box(slide, title, 0.3, 0.2, 9.4, 0.8,
                 font_size=24, bold=True, color=WHITE, align=PP_ALIGN.LEFT)


def add_footer(slide, page_num):
    add_rect(slide, 0, 6.9, 10, 0.6, DARK_GRAY)
    add_text_box(slide, f"Adobe Data Engineer Assessment  |  Search Keyword Performance Analyzer  |  {page_num}",
                 0.3, 6.95, 9.4, 0.4, font_size=9, color=WHITE, align=PP_ALIGN.CENTER)


# ── Slide 1: GitHub Actions CI/CD ─────────────────────────────────────────────
def slide_cicd(prs):
    slide = add_slide(prs)
    clear_slide(slide)
    add_header_bar(slide, "GitHub Actions CI/CD Pipeline")
    add_footer(slide, "CI/CD")

    # Pipeline boxes
    boxes = [
        (0.4,  1.3, 1.8, 0.7, GREEN,      "PUSH / PR\nManual Trigger"),
        (2.4,  1.3, 1.8, 0.7, BLUE,       "Unit Tests\n(python unittest)"),
        (4.4,  1.3, 1.8, 0.7, BLUE,       "Package Lambda\n(zip + artifact)"),
        (6.4,  1.3, 1.8, 0.7, ADOBE_RED,  "Terraform Apply\n(push to main only)"),
        (6.4,  2.4, 1.8, 0.7, DARK_GRAY,  "Terraform Plan\n(PR comment only)"),
    ]
    for l, t, w, h, color, label in boxes:
        r = add_rect(slide, l, t, w, h, color)
        tf = r.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = label
        run.font.size = Pt(9)
        run.font.bold = True
        run.font.color.rgb = WHITE

    # Arrows (simple text)
    for arrow_l, arrow_t in [(2.25, 1.55), (4.25, 1.55)]:
        add_text_box(slide, "▶", arrow_l, arrow_t, 0.2, 0.3,
                     font_size=14, bold=True, color=DARK_GRAY)

    add_text_box(slide, "▶  (push)", 6.25, 1.55, 1.0, 0.3,
                 font_size=9, color=DARK_GRAY)
    add_text_box(slide, "▶  (PR)", 6.25, 2.55, 1.0, 0.3,
                 font_size=9, color=DARK_GRAY)

    # Details section
    details = [
        ("Triggers",        "Push to main  |  Pull Request  |  workflow_dispatch (manual)"),
        ("Test Job",         "Runs unittest — no AWS needed. Catches logic errors early."),
        ("Package Job",      "Zips src/ into lambda.zip. Uploaded as GitHub artifact (30-day retention)."),
        ("Deploy Job",       "terraform apply — only runs on push to main. Deploys Lambda, S3, Glue, Athena."),
        ("Plan Job",         "terraform plan on every PR. Output posted as PR comment for review before merge."),
        ("Branch Protection","Unit Tests + Package Lambda must pass. 1 PR approval required. No direct push to main."),
        ("Secrets",          "AWS_ACCESS_KEY_ID and AWS_SECRET_ACCESS_KEY stored in GitHub Secrets — never in code."),
    ]
    y = 3.3
    for label, value in details:
        add_rect(slide, 0.4, y, 2.0, 0.32, LIGHT_GRAY, DARK_GRAY)
        add_text_box(slide, label, 0.45, y + 0.04, 1.9, 0.28,
                     font_size=9, bold=True, color=DARK_GRAY)
        add_text_box(slide, value, 2.5, y + 0.04, 7.0, 0.28,
                     font_size=9, color=DARK_GRAY)
        y += 0.38


# ── Slide 2: AWS Architecture ──────────────────────────────────────────────────
def slide_architecture(prs):
    slide = add_slide(prs)
    clear_slide(slide)
    add_header_bar(slide, "AWS Architecture — Medallion Lakehouse")
    add_footer(slide, "Architecture")

    layers = [
        (0.3,  1.3, 1.7, 4.8, RGBColor(0xE8,0xF4,0xE8), "TRIGGER",   "S3\nlanding/\ndata.sql"),
        (2.2,  1.3, 1.7, 4.8, RGBColor(0xE8,0xF0,0xF8), "COMPUTE",   "AWS Lambda\nPython 3.12\n512 MB / 5 min"),
        (4.1,  1.3, 1.7, 4.8, RGBColor(0xF8,0xF0,0xE8), "BRONZE",    "S3\nbronze/\nRaw Archive"),
        (6.0,  1.3, 1.7, 4.8, RGBColor(0xF8,0xF8,0xE8), "GOLD",      "S3\ngold/\nOutput .tab"),
        (7.9,  1.3, 1.7, 4.8, RGBColor(0xF0,0xE8,0xF8), "QUERY",     "Athena\n+\nGlue Catalog"),
    ]
    for l, t, w, h, bg, label, content in layers:
        add_rect(slide, l, t, w, h, bg, DARK_GRAY)
        add_text_box(slide, label, l, t + 0.05, w, 0.3,
                     font_size=8, bold=True, color=DARK_GRAY, align=PP_ALIGN.CENTER)
        add_text_box(slide, content, l, t + 0.45, w, h - 0.55,
                     font_size=10, bold=False, color=DARK_GRAY, align=PP_ALIGN.CENTER)

    # Arrows between layers
    for arrow_x in [1.97, 3.87, 5.77, 7.67]:
        add_text_box(slide, "▶", arrow_x, 3.3, 0.25, 0.4,
                     font_size=18, bold=True, color=DARK_GRAY)

    # Supporting services row
    add_text_box(slide, "Supporting Services", 0.3, 6.1, 9.5, 0.3,
                 font_size=10, bold=True, color=DARK_GRAY)
    services = ["KMS Encryption", "CloudWatch Logs", "IAM Least Privilege",
                "S3 Lifecycle Rules", "VPC-ready IAM"]
    x = 0.3
    for svc in services:
        add_rect(slide, x, 6.45, 1.75, 0.35, DARK_GRAY)
        add_text_box(slide, svc, x + 0.05, 6.48, 1.65, 0.28,
                     font_size=8, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        x += 1.9


# ── Slide 3: Security & Encryption ────────────────────────────────────────────
def slide_security(prs):
    slide = add_slide(prs)
    clear_slide(slide)
    add_header_bar(slide, "Security & Encryption")
    add_footer(slide, "Security")

    items = [
        ("S3 Encryption",           "SSE-KMS with customer-managed key. Bucket key enabled — reduces KMS API calls by ~99%."),
        ("Glue Data Catalog",       "Catalog metadata encrypted with same KMS key. Connection passwords encrypted at rest."),
        ("Athena Query Results",    "Output stored in S3 with SSE-KMS encryption. 100 MB per-query scan limit to control costs."),
        ("KMS Key Rotation",        "Automatic annual rotation enabled. 7-day deletion window prevents accidental data loss."),
        ("S3 Public Access",        "All public access fully blocked: ACLs, policies, public buckets — 4/4 flags enabled."),
        ("IAM Least Privilege",     "Lambda role has only: s3:Get/Put/Copy on bucket, kms:Decrypt/GenerateDataKey, logs:*."),
        ("GitHub Secrets",          "AWS credentials stored as encrypted GitHub Secrets. Never hardcoded in source or config."),
        ("Branch Protection",       "Direct push to main blocked. PRs require 1 approval + passing CI before merge."),
        ("No External Dependencies","Python standard library only. Zero third-party packages = zero supply chain risk."),
    ]

    y = 1.25
    for i, (title, detail) in enumerate(items):
        bg = LIGHT_GRAY if i % 2 == 0 else WHITE
        add_rect(slide, 0.3, y, 9.4, 0.52, bg)
        add_text_box(slide, f"✓  {title}", 0.4, y + 0.04, 2.5, 0.44,
                     font_size=9, bold=True, color=ADOBE_RED)
        add_text_box(slide, detail, 3.0, y + 0.08, 6.6, 0.4,
                     font_size=9, color=DARK_GRAY)
        y += 0.56


# ── Slide 4: Branch Protection & Git Workflow ─────────────────────────────────
def slide_git(prs):
    slide = add_slide(prs)
    clear_slide(slide)
    add_header_bar(slide, "Branch Protection & Git Workflow")
    add_footer(slide, "Git Workflow")

    # Left column — workflow steps
    add_text_box(slide, "Developer Workflow", 0.3, 1.25, 4.5, 0.35,
                 font_size=13, bold=True, color=DARK_GRAY)
    steps = [
        "1.  git checkout main && git pull",
        "2.  git checkout -b feature/SKA-123-description",
        "3.  Make changes & commit (Conventional Commits)",
        "4.  git push -u origin feature/SKA-123-description",
        "5.  gh pr create  (opens Pull Request)",
        "6.  CI runs: Unit Tests + Package Lambda",
        "7.  Reviewer approves PR",
        "8.  Merge via GitHub UI (Squash and merge)",
        "9.  Delete feature branch",
    ]
    y = 1.7
    for step in steps:
        add_text_box(slide, step, 0.4, y, 4.4, 0.35, font_size=9, color=DARK_GRAY)
        y += 0.38

    # Right column — branch naming
    add_text_box(slide, "Branch Naming Standards", 5.1, 1.25, 4.5, 0.35,
                 font_size=13, bold=True, color=DARK_GRAY)
    naming = [
        ("feature/", "feature/SKA-42-yahoo-parser"),
        ("fix/",     "fix/SKA-55-revenue-rounding"),
        ("hotfix/",  "hotfix/SKA-99-lambda-oom"),
        ("release/", "release/1.2.0"),
        ("chore/",   "chore/update-dependencies"),
    ]
    y = 1.7
    for prefix, example in naming:
        add_rect(slide, 5.1, y, 1.1, 0.32, ADOBE_RED)
        add_text_box(slide, prefix, 5.12, y + 0.04, 1.06, 0.26,
                     font_size=9, bold=True, color=WHITE)
        add_text_box(slide, example, 6.3, y + 0.04, 2.8, 0.28,
                     font_size=9, color=DARK_GRAY)
        y += 0.42

    add_text_box(slide, "Commit Format (Conventional Commits)", 5.1, 4.0, 4.5, 0.35,
                 font_size=11, bold=True, color=DARK_GRAY)
    add_text_box(slide,
                 "feat(scope): description\nfix(scope): description\n"
                 "docs / test / refactor / ci / chore",
                 5.1, 4.4, 4.5, 0.8, font_size=9, color=DARK_GRAY)

    # Protection rules box
    add_rect(slide, 0.3, 5.4, 9.4, 1.2, LIGHT_GRAY, DARK_GRAY)
    add_text_box(slide, "Branch Protection Rules — main", 0.5, 5.45, 5.0, 0.35,
                 font_size=11, bold=True, color=ADOBE_RED)
    rules = ("✓  Direct push to main blocked     ✓  1 PR approval required\n"
             "✓  Unit Tests must pass             ✓  Package Lambda must pass\n"
             "✓  Force push disabled              ✓  Branch deletion disabled")
    add_text_box(slide, rules, 0.5, 5.85, 9.0, 0.7, font_size=9, color=DARK_GRAY)


# ── Slide 5: Scalability ───────────────────────────────────────────────────────
def slide_scalability(prs):
    slide = add_slide(prs)
    clear_slide(slide)
    add_header_bar(slide, "Scalability — Handling 10 GB+ Files")
    add_footer(slide, "Scalability")

    add_text_box(slide,
                 "Current design streams the file row by row (csv.DictReader) — "
                 "memory usage is O(unique visitors), not O(file size). "
                 "This handles moderate files efficiently but has limits at 10 GB+ scale.",
                 0.3, 1.2, 9.4, 0.65, font_size=10, color=DARK_GRAY)

    headers = ["Approach", "File Size", "How", "AWS Service"]
    col_widths = [1.8, 1.2, 4.2, 1.8]
    col_x = [0.3, 2.1, 3.3, 7.5]

    # Header row
    x = 0.3
    for i, h in enumerate(headers):
        add_rect(slide, col_x[i], 2.0, col_widths[i], 0.35, DARK_GRAY)
        add_text_box(slide, h, col_x[i] + 0.05, 2.03, col_widths[i] - 0.1, 0.3,
                     font_size=9, bold=True, color=WHITE)

    rows = [
        ("Chunked multiprocessing", "10–50 GB", "Split file into N chunks, process in parallel, merge results", "EC2 / Lambda"),
        ("AWS Glue (PySpark)",      "50+ GB",   "Serverless managed Spark. Reads from S3, partitioned automatically", "Glue"),
        ("Amazon EMR",              "100+ GB",  "Full Spark cluster. Maximum control, optimal for recurring large jobs", "EMR"),
        ("Kinesis Streams",         "Real-time","Process each hit as it arrives. No batch window — sub-second latency", "Kinesis"),
    ]
    y = 2.45
    for i, row in enumerate(rows):
        bg = LIGHT_GRAY if i % 2 == 0 else WHITE
        add_rect(slide, 0.3, y, 9.4, 0.45, bg)
        for j, (cell, cw, cx) in enumerate(zip(row, col_widths, col_x)):
            add_text_box(slide, cell, cx + 0.05, y + 0.06, cw - 0.1, 0.35,
                         font_size=9, color=DARK_GRAY)
        y += 0.48

    add_text_box(slide, "Current Bottlenecks at Scale", 0.3, 4.7, 5.0, 0.35,
                 font_size=11, bold=True, color=ADOBE_RED)
    bottlenecks = [
        "Single-threaded processing — one core for entire file",
        "In-memory visitor map — grows with unique IPs (millions at scale)",
        "Lambda limits — 15-min timeout, 10 GB /tmp storage",
        "Single-file input — no native partitioning or parallelism",
    ]
    y = 5.1
    for b in bottlenecks:
        add_text_box(slide, f"•  {b}", 0.4, y, 9.0, 0.32, font_size=9, color=DARK_GRAY)
        y += 0.35


# ── Main ───────────────────────────────────────────────────────────────────────
prs = Presentation(INPUT)

slide_cicd(prs)
slide_architecture(prs)
slide_security(prs)
slide_git(prs)
slide_scalability(prs)

prs.save(OUTPUT)
print(f"Saved {OUTPUT} with {len(prs.slides)} slides total.")
